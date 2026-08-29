# Visor de diapositivas — sustituto de LibreOffice, que en este contenedor no
# convierte ni un archivo de texto.
#
# NO dibuja lo que yo creía haber puesto: lee el .pptx YA GENERADO con
# python-pptx y dibuja lo que de verdad quedó adentro (posición, tamaño,
# color, fuente y texto de cada forma). Por eso sirve como revisión de verdad
# y no como espejo de mis intenciones.
#
#     python3 vista.py sistema-ciudad-maderas.pptx
import base64, html, sys, os
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU_IN = 914400
PX = 100  # píxeles por pulgada

ALIGN = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}
ANCHOR = {MSO_ANCHOR.TOP: "flex-start", MSO_ANCHOR.MIDDLE: "center", MSO_ANCHOR.BOTTOM: "flex-end"}


def color_de(fmt):
    """Color sólido de un relleno o línea, o None."""
    try:
        if fmt.type is None:
            return None
        c = fmt.fore_color
        if c.type is not None and hasattr(c, "rgb"):
            return "#" + str(c.rgb)
    except Exception:
        pass
    return None


A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def transparencia(shape):
    """Alfa del RELLENO — nunca la de la sombra.

    Buscar `.//alpha` a ciegas encontraba primero el alfa del efecto de sombra
    y pintaba translúcida una tarjeta opaca: los recuadros azul marino salían
    grises y parecían un error de la presentación que no existía.
    """
    try:
        spPr = shape._element.find(f".//{A}prstGeom/..")
        if spPr is None:
            return 1.0
        relleno = spPr.find(f"{A}solidFill")
        if relleno is None:
            return 1.0
        alpha = relleno.find(f".//{A}alpha")
        if alpha is not None:
            return int(alpha.get("val")) / 100000.0
    except Exception:
        pass
    return 1.0


def geometria(shape):
    """El nombre de la forma preestablecida: rect, roundRect, ellipse…"""
    try:
        g = shape._element.find(f".//{A}prstGeom")
        if g is not None:
            return g.get("prst") or "rect"
    except Exception:
        pass
    return "rect"


def redondeo(shape):
    g = geometria(shape)
    if g == "ellipse":
        return "50%"
    if g == "roundRect":
        return "10px"
    return "0"


def render(ruta, salida="vista"):
    pr = Presentation(ruta)
    W = pr.slide_width / EMU_IN * PX
    H = pr.slide_height / EMU_IN * PX
    partes = []

    for n, sl in enumerate(pr.slides, 1):
        piezas = []
        fondo = "#FFFFFF"
        try:
            c = color_de(sl.background.fill)
            if c:
                fondo = c
        except Exception:
            pass

        for sh in sl.shapes:
            x = sh.left / EMU_IN * PX if sh.left is not None else 0
            y = sh.top / EMU_IN * PX if sh.top is not None else 0
            w = sh.width / EMU_IN * PX if sh.width is not None else 0
            h = sh.height / EMU_IN * PX if sh.height is not None else 0
            base = f"position:absolute;left:{x:.1f}px;top:{y:.1f}px;width:{w:.1f}px;height:{h:.1f}px;"

            # Imágenes
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                img = sh.image
                b64 = base64.b64encode(img.blob).decode()
                piezas.append(
                    f'<img style="{base}object-fit:cover" src="data:{img.content_type};base64,{b64}">')
                continue

            # Relleno y borde
            estilo = base
            try:
                cf = color_de(sh.fill)
                if cf:
                    a = transparencia(sh)
                    if a < 1:
                        r, g, b = int(cf[1:3], 16), int(cf[3:5], 16), int(cf[5:7], 16)
                        estilo += f"background:rgba({r},{g},{b},{a:.2f});"
                    else:
                        estilo += f"background:{cf};"
            except Exception:
                pass
            try:
                cl = color_de(sh.line.fill) if sh.line.fill.type is not None else None
                if cl:
                    ancho = (sh.line.width / EMU_IN * PX) if sh.line.width else 1
                    estilo += f"border:{max(1, ancho):.1f}px solid {cl};box-sizing:border-box;"
            except Exception:
                pass
            estilo += f"border-radius:{redondeo(sh)};"

            # Texto
            interior = ""
            if sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                anc = ANCHOR.get(tf.vertical_anchor, "flex-start")
                parrafos = []
                for para in tf.paragraphs:
                    al = ALIGN.get(para.alignment, "left")
                    runs = []
                    for r in para.runs:
                        f = r.font
                        st = ""
                        if f.size:
                            st += f"font-size:{f.size.pt * PX / 72:.1f}px;"
                        if f.name:
                            st += f"font-family:'{f.name}',serif;"
                        if f.bold:
                            st += "font-weight:700;"
                        if f.italic:
                            st += "font-style:italic;"
                        try:
                            if f.color and f.color.type is not None:
                                st += f"color:#{f.color.rgb};"
                        except Exception:
                            pass
                        runs.append(f'<span style="{st}">{html.escape(r.text)}</span>')
                    if not runs:
                        runs = ["&nbsp;"]
                    parrafos.append(f'<div style="text-align:{al}">{"".join(runs)}</div>')
                interior = (
                    f'<div style="position:absolute;inset:0;display:flex;flex-direction:column;'
                    f'justify-content:{anc};overflow:visible">{"".join(parrafos)}</div>')

            piezas.append(f'<div style="{estilo}">{interior}</div>')

        partes.append(
            f'<div class="lam"><div class="num">{n}</div>'
            f'<div class="d" style="background:{fondo}">{"".join(piezas)}</div></div>')

    doc = f"""<!doctype html><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#2b3138;padding:26px;font-family:Calibri,sans-serif}}
.lam{{margin-bottom:30px}}
.num{{color:#9fb0bd;font:600 15px Calibri;margin-bottom:7px}}
.d{{position:relative;width:{W:.0f}px;height:{H:.0f}px;overflow:hidden;
    box-shadow:0 6px 26px rgba(0,0,0,.5)}}
</style>{"".join(partes)}"""
    open(f"{salida}.html", "w", encoding="utf-8").write(doc)
    print(f"{salida}.html — {len(pr.slides)} diapositivas de {W:.0f}x{H:.0f}px")
    return W, H, len(pr.slides)


if __name__ == "__main__":
    render(sys.argv[1] if len(sys.argv) > 1 else "sistema-ciudad-maderas.pptx")
